#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""verify.py — 验证 PPTX 内容完整性"""
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

# 强制 UTF-8 输出
if sys.platform == 'win32':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Pt

def inspect_pptx(path):
    prs = Presentation(path)
    print(f"幻灯片数: {len(prs.slides)}")
    print(f"尺寸: {prs.slide_width} x {prs.slide_height} EMU")
    print()

    for i, slide in enumerate(prs.slides):
        print(f"══ 第 {i+1} 页 ══")
        print(f"  形状数: {len(slide.shapes)}")
        for shape in slide.shapes:
            print(f"  ├── {shape.shape_type} | name={shape.name}")
            if shape.has_text_frame:
                tf = shape.text_frame
                for j, para in enumerate(tf.paragraphs):
                    text = ''.join(r.text for r in para.runs)
                    if text:
                        font_info = ""
                        if para.runs:
                            r = para.runs[0]
                            font_info = f" [size={r.font.size}, bold={r.font.bold}, italic={r.font.italic}]"
                        print(f"  │   段落{j}: \"{text[:80]}\"{font_info}")
            elif shape.has_table:
                tbl = shape.table
                print(f"  │   表格: {len(tbl.rows)}行 x {len(tbl.columns)}列")
                for r_idx, row in enumerate(tbl.rows):
                    cells_text = []
                    for cell in row.cells:
                        cell_text = cell.text_frame.text[:20]
                        cells_text.append(cell_text)
                    print(f"  │   行{r_idx}: {cells_text}")
            elif shape.shape_type == 13:  # PICTURE
                print(f"  │   图片")
        print()

if __name__ == '__main__':
    path = sys.argv[1] if len(sys.argv) > 1 else os.path.join(os.path.dirname(__file__), 'demo_output.pptx')
    inspect_pptx(path)
