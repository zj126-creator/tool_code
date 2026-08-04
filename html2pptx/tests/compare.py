#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""compare.py — 对比 HTML 原文与 PPTX 内容，确保一模一样"""
import sys, os, io, re
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from bs4 import BeautifulSoup
from pptx import Presentation

def extract_html_text(html_path):
    with open(html_path, 'r', encoding='utf-8') as f:
        content = f.read()
    soup = BeautifulSoup(content, 'lxml')
    for tag in soup.find_all(['script','style','noscript']):
        tag.decompose()
    # 提取所有文本
    texts = []
    for elem in soup.find_all(string=True):
        text = str(elem).strip()
        if text:
            texts.append(text)
    return texts

def extract_pptx_text(pptx_path):
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = ''.join(r.text for r in para.runs).strip()
                    if text:
                        texts.append(text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text_frame.text.strip()
                        if text:
                            texts.append(text)
    return texts

def compare(html_path, pptx_path):
    html_texts = extract_html_text(html_path)
    pptx_texts = extract_pptx_text(pptx_path)

    print(f"HTML 文本片段数: {len(html_texts)}")
    print(f"PPTX 文本片段数: {len(pptx_texts)}")
    print()

    # 检查每个 HTML 文本片段是否在 PPTX 中存在
    missing = []
    for ht in html_texts:
        found = False
        for pt in pptx_texts:
            if ht in pt or pt in ht:
                found = True
                break
        if not found:
            missing.append(ht)

    if missing:
        print(f"⚠️ 有 {len(missing)} 个文本片段未在 PPTX 中找到:")
        for m in missing:
            print(f"  - \"{m[:60]}\"")
    else:
        print("✅ 所有 HTML 文本片段都在 PPTX 中找到!")

    # 反向检查
    extra = []
    for pt in pptx_texts:
        found = False
        for ht in html_texts:
            if pt in ht or ht in pt:
                found = True
                break
        if not found:
            extra.append(pt)

    if extra:
        print(f"\n⚠️ PPTX 中有 {len(extra)} 个额外文本片段（可能来自格式拆分）:")
        for e in extra[:10]:
            print(f"  + \"{e[:60]}\"")
    else:
        print("✅ PPTX 中没有多余内容!")

    return len(missing) == 0

if __name__ == '__main__':
    html_path = sys.argv[1] if len(sys.argv) > 1 else 'demo.html'
    pptx_path = sys.argv[2] if len(sys.argv) > 2 else 'demo_output.pptx'
    ok = compare(html_path, pptx_path)
    sys.exit(0 if ok else 1)
