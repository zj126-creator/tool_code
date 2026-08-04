#!/usr/bin/env python3
"""
pptx_builder.py — 将 SlideContent 列表写入可编辑的 PPTX 文件

第一性原理: PPTX 格式的每个元素都是 XML 描述的原生对象。
构建器的职责是将结构化数据映射为 PPTX 原生元素（文本框、图片、表格），
确保内容完整、格式正确、元素可编辑。
"""
import io
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from models import ContentBlock, SlideContent, Paragraph, TextRun


# PPTX 字号有效范围 (EMU)
MIN_FONT_EMU = 100   # 1pt
MAX_FONT_EMU = 400000  # 4000pt


class PPTXBuilder:
    """将 SlideContent 列表构建为 PPTX 文件"""

    DEFAULT_FONT_SIZE = 14
    DEFAULT_TITLE_SIZES = {1: 32, 2: 28, 3: 24, 4: 20, 5: 18, 6: 16}
    DEFAULT_FONT_NAME = '微软雅黑'

    def __init__(self, slide_width_in: float = 10.0, slide_height_in: float = 7.5):
        self.slide_width = slide_width_in
        self.slide_height = slide_height_in
        self.margin = 0.5
        self.content_width = slide_width_in - 2 * self.margin
        self.content_height = slide_height_in - 2 * self.margin

    def build(self, slides: List[SlideContent], output_path: str):
        """构建 PPTX 文件"""
        prs = Presentation()
        prs.slide_width = Inches(self.slide_width)
        prs.slide_height = Inches(self.slide_height)
        blank_layout = prs.slide_layouts[6]

        for slide_content in slides:
            slide = prs.slides.add_slide(blank_layout)
            self._render_slide(slide, slide_content, prs)

        prs.save(output_path)

    def _render_slide(self, slide, slide_content: SlideContent, prs):
        """渲染一页 slide"""
        y = self.margin
        x = self.margin
        width = self.content_width

        for block in slide_content.blocks:
            if block.type == 'spacer':
                y += 0.15
                continue
            if block.type == 'divider':
                self._add_divider(slide, x, y, width)
                y += 0.2
                continue
            if block.type == 'title':
                y = self._add_title(slide, block, x, y, width)
                continue
            if block.type == 'paragraph':
                y = self._add_paragraph_block(slide, block, x, y, width)
                continue
            if block.type == 'image':
                y = self._add_image(slide, block, x, y, width)
                continue
            if block.type == 'code_block':
                y = self._add_code_block(slide, block, x, y, width)
                continue
            if block.type == 'table':
                y = self._add_table(slide, block, x, y, width)
                continue

    def _add_divider(self, slide, x, y, width):
        """添加分隔线"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(width), Pt(1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.fill.background()

    def _add_title(self, slide, block: ContentBlock, x, y, width) -> float:
        """添加标题，返回新的 y"""
        font_size = self.DEFAULT_TITLE_SIZES.get(block.level, 16)
        height_in = font_size / 72.0 * 1.5

        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(width), Inches(height_in)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.TOP

        first = True
        for para in block.paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = self._get_alignment(para.alignment)
            p.space_before = Pt(0)
            p.space_after = Pt(4)

            for run in para.runs:
                run_obj = p.add_run()
                run_obj.text = run.text
                run_obj.font.size = self._safe_pt(run.font_size or font_size)
                run_obj.font.bold = run.bold or True
                run_obj.font.italic = run.italic
                run_obj.font.underline = run.underline
                run_obj.font.name = run.font_name or self.DEFAULT_FONT_NAME
                if run.color and run.color.startswith('#'):
                    run_obj.font.color.rgb = self._hex_to_rgb(run.color)
                else:
                    run_obj.font.color.rgb = RGBColor(0x1a, 0x1a, 0x1a)

        return y + height_in + 0.15

    def _add_paragraph_block(self, slide, block: ContentBlock, x, y, width) -> float:
        """添加段落块，返回新的 y"""
        total_text = ''
        for para in block.paragraphs:
            for run in para.runs:
                total_text += run.text
            total_text += '\n'
        font_size = self._get_block_font_size(block) or self.DEFAULT_FONT_SIZE
        # 安全字号
        font_size = max(1.0, font_size)
        line_height = font_size / 72.0 * 1.5
        chars_per_line = max(20, int(width * 72 / font_size * 1.8))
        total_lines = 0
        for para in block.paragraphs:
            para_text = ''.join(r.text for r in para.runs)
            if not para_text:
                total_lines += 1
            else:
                # 处理文本中的换行
                for line in para_text.split('\n'):
                    if not line:
                        total_lines += 1
                    else:
                        total_lines += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
        height_in = max(0.3, total_lines * line_height + 0.1)

        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(width), Inches(height_in)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = MSO_ANCHOR.TOP

        if block.bg_color:
            self._set_textbox_bg(txBox, block.bg_color)

        first = True
        for para in block.paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = self._get_alignment(para.alignment)
            p.space_before = Pt(para.space_before)
            p.space_after = Pt(para.space_after)
            if para.line_spacing:
                p.line_spacing = para.line_spacing

            if para.bullet:
                self._add_bullet(p, para.bullet_level,
                                ordered=para.ordered,
                                list_number=para.list_number)

            for run in para.runs:
                run_obj = p.add_run()
                run_obj.text = run.text
                run_obj.font.size = self._safe_pt(run.font_size or font_size)
                run_obj.font.bold = run.bold
                run_obj.font.italic = run.italic
                run_obj.font.underline = run.underline
                run_obj.font.name = run.font_name or self.DEFAULT_FONT_NAME
                if run.color and run.color.startswith('#'):
                    run_obj.font.color.rgb = self._hex_to_rgb(run.color)
                else:
                    run_obj.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                # 删除线
                if run.strikethrough:
                    self._set_strikethrough(run_obj)

        return y + height_in

    def _add_image(self, slide, block: ContentBlock, x, y, width) -> float:
        """添加图片，返回新的 y"""
        if not block.image_data:
            return y + 0.2
        try:
            if block.image_ext in ('png','jpg','jpeg','gif','bmp'):
                img_stream = io.BytesIO(block.image_data)
                try:
                    from PIL import Image
                    img = Image.open(img_stream)
                    iw, ih = img.size
                    img_stream.seek(0)
                except Exception:
                    iw, ih = 400, 300
                max_w = width
                max_h = 4.0
                ratio = min(max_w / (iw / 96.0), max_h / (ih / 96.0)) if iw > 0 and ih > 0 else 1
                w_in = (iw / 96.0) * ratio
                h_in = (ih / 96.0) * ratio
                x_in = x + (width - w_in) / 2
                slide.shapes.add_picture(img_stream, Inches(x_in), Inches(y),
                                         Inches(w_in), Inches(h_in))
                new_y = y + h_in + 0.1
            elif block.image_ext == 'svg':
                try:
                    import cairosvg
                    png_data = cairosvg.svg2png(block.image_data, output_width=800)
                    img_stream = io.BytesIO(png_data)
                    slide.shapes.add_picture(img_stream, Inches(x), Inches(y),
                                             Inches(min(width, 6)), Inches(3))
                    new_y = y + 3.1
                except ImportError:
                    alt_text = block.image_alt or '[SVG 图片]'
                    new_y = self._add_image_placeholder(slide, x, y, width, alt_text)
            elif block.image_ext == 'webp':
                try:
                    from PIL import Image
                    img = Image.open(io.BytesIO(block.image_data))
                    if img.mode in ('RGBA','P'):
                        img = img.convert('RGBA')
                    png_stream = io.BytesIO()
                    img.save(png_stream, format='PNG')
                    png_stream.seek(0)
                    iw, ih = img.size
                    max_w = width
                    max_h = 4.0
                    ratio = min(max_w / (iw/96.0), max_h / (ih/96.0)) if iw > 0 and ih > 0 else 1
                    w_in = (iw/96.0) * ratio
                    h_in = (ih/96.0) * ratio
                    x_in = x + (width - w_in) / 2
                    slide.shapes.add_picture(png_stream, Inches(x_in), Inches(y),
                                             Inches(w_in), Inches(h_in))
                    new_y = y + h_in + 0.1
                except Exception:
                    alt_text = block.image_alt or '[WebP 图片]'
                    new_y = self._add_image_placeholder(slide, x, y, width, alt_text)
            else:
                alt_text = block.image_alt or f'[{block.image_ext} 图片]'
                new_y = self._add_image_placeholder(slide, x, y, width, alt_text)
        except Exception as e:
            alt_text = block.image_alt or '[图片]'
            new_y = self._add_image_placeholder(slide, x, y, width, f'{alt_text} (加载失败: {e})')

        # 图片说明
        if block.image_alt and new_y < self.slide_height - self.margin - 0.3:
            cap_box = slide.shapes.add_textbox(
                Inches(x), Inches(new_y), Inches(width), Inches(0.3)
            )
            tf = cap_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = block.image_alt
            run.font.size = Pt(10)
            run.font.italic = True
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            new_y += 0.35

        return new_y

    def _add_image_placeholder(self, slide, x, y, width, text) -> float:
        """添加图片占位文本"""
        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(width), Inches(0.5)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f'[{text}]'
        run.font.italic = True
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        return y + 0.6

    def _add_code_block(self, slide, block: ContentBlock, x, y, width) -> float:
        """添加代码块"""
        lines = block.code_text.split('\n')
        line_height = 0.2
        height_in = max(0.4, len(lines) * line_height + 0.1)

        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(width), Inches(height_in)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        self._set_textbox_bg(txBox, '#1E1E1E')

        first = True
        for line in lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0
            run = p.add_run()
            run.text = line
            run.font.size = Pt(10)
            run.font.name = 'Consolas'
            run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)

        return y + height_in + 0.1

    def _add_table(self, slide, block: ContentBlock, x, y, width) -> float:
        """添加表格"""
        if not block.table_data:
            return y + 0.2
        rows = len(block.table_data)
        cols = len(block.table_data[0]) if block.table_data else 0
        if cols == 0:
            return y + 0.2

        row_height = 0.3
        height_in = rows * row_height
        col_width = width / cols

        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(x), Inches(y),
            Inches(width), Inches(height_in)
        )
        table = table_shape.table

        for i in range(cols):
            table.columns[i].width = Inches(col_width)

        # 填充数据
        for r_idx, row_data in enumerate(block.table_data):
            for c_idx, cell_paras in enumerate(row_data):
                if c_idx >= cols:
                    break
                cell = table.cell(r_idx, c_idx)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = cell.text_frame
                tf.word_wrap = True
                first = True
                for para in cell_paras:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.alignment = self._get_alignment(para.alignment)
                    p.space_before = Pt(2)
                    p.space_after = Pt(2)
                    for run in para.runs:
                        run_obj = p.add_run()
                        run_obj.text = run.text
                        run_obj.font.size = self._safe_pt(run.font_size or 11)
                        run_obj.font.bold = run.bold
                        run_obj.font.italic = run.italic
                        run_obj.font.underline = run.underline
                        run_obj.font.name = run.font_name or self.DEFAULT_FONT_NAME
                        if run.color and run.color.startswith('#'):
                            run_obj.font.color.rgb = self._hex_to_rgb(run.color)
                        else:
                            run_obj.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        if run.strikethrough:
                            self._set_strikethrough(run_obj)

                # 表头样式
                is_header = block.table_header and r_idx == 0
                if is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
                    for p in tf.paragraphs:
                        for r in p.runs:
                            r.font.bold = True
                            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    if r_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 处理合并单元格
        if block.table_merges:
            for merge_info in block.table_merges:
                if len(merge_info) == 4:
                    start_row, start_col, rowspan, colspan = merge_info
                    end_row = start_row + rowspan - 1
                    end_col = start_col + colspan - 1
                    try:
                        # python-pptx 的 merge_cells 方法
                        cell_start = table.cell(start_row, start_col)
                        cell_end = table.cell(end_row, end_col)
                        cell_start.merge(cell_end)
                    except Exception:
                        pass  # 合并失败不影响内容

        return y + height_in + 0.15

    def _set_textbox_bg(self, textbox, color_hex: str):
        """设置文本框背景色"""
        rgb = self._hex_to_rgb(color_hex)
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = rgb

    def _get_alignment(self, align_str: Optional[str]):
        """获取对齐方式"""
        if not align_str:
            return PP_ALIGN.LEFT
        mapping = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY,
        }
        return mapping.get(align_str.lower(), PP_ALIGN.LEFT)

    def _hex_to_rgb(self, hex_str: str) -> RGBColor:
        """hex 字符串 → RGBColor"""
        hex_str = hex_str.strip().lstrip('#')
        if len(hex_str) == 3:
            hex_str = ''.join(c*2 for c in hex_str)
        if len(hex_str) >= 6:
            try:
                r = int(hex_str[0:2], 16)
                g = int(hex_str[2:4], 16)
                b = int(hex_str[4:6], 16)
                return RGBColor(r, g, b)
            except ValueError:
                pass
        return RGBColor(0x33, 0x33, 0x33)

    def _get_block_font_size(self, block: ContentBlock) -> Optional[float]:
        """获取块的默认字号"""
        for para in block.paragraphs:
            for run in para.runs:
                if run.font_size:
                    return run.font_size
        return None

    def _safe_pt(self, size_pt) -> Pt:
        """将字号转换为安全的 Pt 对象，确保在 PPTX 合法范围内"""
        try:
            size_pt = float(size_pt)
        except (TypeError, ValueError):
            size_pt = self.DEFAULT_FONT_SIZE
        # PPTX 字号范围: 1pt ~ 4000pt
        size_pt = max(1.0, min(4000.0, size_pt))
        return Pt(size_pt)

    def _set_strikethrough(self, run_obj):
        """给 run 设置删除线"""
        rPr = run_obj._r.get_or_add_rPr()
        # 添加 strike 属性
        rPr.set('strike', 'sngStrike')

    def _add_bullet(self, paragraph, level: int, ordered: bool = False, list_number: int = 0):
        """给段落添加项目符号"""
        pPr = paragraph._p.get_or_add_pPr()

        marL = 274320 * (level + 1)
        indent = -228600
        pPr.set('marL', str(marL))
        pPr.set('indent', str(indent))

        # 移除已有的项目符号元素
        for child in list(pPr):
            tag = child.tag
            if tag in (qn('a:buNone'), qn('a:buChar'), qn('a:buAutoNum'),
                      qn('a:buFont'), qn('a:buClr'), qn('a:buSzPct')):
                pPr.remove(child)

        # buFont 必须在 buChar/buAutoNum 之前
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
        pPr.append(buFont)

        if ordered and list_number > 0:
            # 有序列表 → buAutoNum
            buAutoNum = pPr.makeelement(qn('a:buAutoNum'), {'type': 'arabicPeriod'})
            pPr.append(buAutoNum)
        else:
            # 无序列表 → buChar
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
            pPr.append(buChar)
